import telebot
import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from deep_translator import GoogleTranslator
import time

# YANGI TOKEN
TOKEN = '8267155928:AAGSfZpZs_ibOBnpFRnd4IoETLJida7-GmI'
bot = telebot.TeleBot(TOKEN)

def get_image(query):
    try:
        translated = GoogleTranslator(source='auto', target='en').translate(query)
        search_word = translated.split()[-1]
        url = f"https://source.unsplash.com/1200x800/?{search_word}"
        res = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'}, timeout=15)
        if res.status_code == 200:
            path = f"img_{int(time.time())}.jpg"
            with open(path, "wb") as f: f.write(res.content)
            return path
    except: return None

@bot.message_handler(commands=['start'])
def start(message):
    markup = telebot.types.ReplyKeyboardMarkup(resize_keyboard=True)
    markup.add(telebot.types.KeyboardButton("📊 Slayt tayyorlash"))
    bot.send_message(message.chat.id, "🌟 **Gamma AI Slayd Bot!**\nTugmani bosing va mavzuni yuboring.", reply_markup=markup)

@bot.message_handler(func=lambda m: m.text == "📊 Slayt tayyorlash")
def ask_topic(message):
    msg = bot.send_message(message.chat.id, "📝 **Slayd mavzusini kiriting:**")
    bot.register_next_step_handler(msg, create_ppt)

def create_ppt(message):
    mavzu = message.text
    bot.send_message(message.chat.id, f"🚀 '{mavzu}' bo'yicha 15 betlik slayd tayyorlanmoqda...")
    
    try:
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
        reja = ["Kirish", "Dolzarblik", "Maqsad", "Tarix", "Tahlil", "Statistika", "Muammolar", "Yechimlar", "Tajriba", "Amaliyot", "Innovatsiya", "Kelajak", "Xulosa", "Takliflar", "Manbalar"]

        for i, qism in enumerate(reja):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(10, 15, 30)
            
            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            p = title.text_frame.paragraphs[0]
            p.text = f"{i+1}. {qism}"
            p.font.bold, p.font.size, p.font.color.rgb = True, Pt(38), RGBColor(0, 255, 180)

            img = get_image(f"{mavzu} {qism}")
            if i % 2 == 0:
                if img: slide.shapes.add_picture(img, Inches(7.2), Inches(1.3), Inches(5.8), Inches(5.5))
                body_pos = (Inches(0.5), Inches(1.8), Inches(6.2), Inches(5))
            else:
                if img: slide.shapes.add_picture(img, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5))
                body_pos = (Inches(6.8), Inches(1.8), Inches(6.2), Inches(5))

            body = slide.shapes.add_textbox(*body_pos)
            tf = body.text_frame
            tf.word_wrap = True
            cp = tf.paragraphs[0]
            cp.text = f"{mavzu} haqida {qism.lower()} tahlili. Gamma AI premium dizayni."
            cp.font.size, cp.font.color.rgb = Pt(24), RGBColor(255, 255, 255)
            if img and os.path.exists(img): os.remove(img)

        name = f"Gamma_{message.chat.id}.pptx"
        prs.save(name)
        with open(name, 'rb') as f: bot.send_document(message.chat.id, f)
        os.remove(name)
    except Exception as e:
        bot.send_message(message.chat.id, f"❌ Xato: {str(e)}")

print("Bot ishga tushdi...")
bot.infinity_polling()